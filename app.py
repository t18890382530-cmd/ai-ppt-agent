import io
import json
import os
import re
from typing import Any

import streamlit as st
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

load_dotenv()

api_key = os.getenv("DEEPSEEK_API_KEY")

if not api_key:
    try:
        api_key = st.secrets["DEEPSEEK_API_KEY"]
    except Exception:
        api_key = None

st.set_page_config(page_title="AI PPT Agent Pro", page_icon="✨", layout="wide")

THEMES = {
    "曜石蓝 Pro": {
        "bg": "F5F7FB",
        "ink": "102A43",
        "body": "334E68",
        "muted": "7B91A7",
        "primary": "153E75",
        "accent": "2F80ED",
        "accent2": "56CCF2",
        "soft": "EAF2FB",
        "card": "FFFFFF",
        "dark": "081B33",
        "white": "FFFFFF",
    },
    "青绿科技 Pro": {
        "bg": "F3FAF8",
        "ink": "123C3A",
        "body": "315B57",
        "muted": "6C918C",
        "primary": "0E6259",
        "accent": "27AE8F",
        "accent2": "6EE7C8",
        "soft": "DDF5EE",
        "card": "FFFFFF",
        "dark": "092E2C",
        "white": "FFFFFF",
    },
    "暖橙教育 Pro": {
        "bg": "FFF8F1",
        "ink": "5A321E",
        "body": "6F4E3D",
        "muted": "A9836C",
        "primary": "A64B1D",
        "accent": "F2994A",
        "accent2": "FFD19A",
        "soft": "FDEBDD",
        "card": "FFFFFF",
        "dark": "3D2114",
        "white": "FFFFFF",
    },
}


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))


def fill(shape, color: str, transparency: int = 0, line: bool = False) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency
    if not line:
        shape.line.fill.background()


def text(
    slide,
    x,
    y,
    w,
    h,
    value: str,
    size: int,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    r = p.runs[0]
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def safe_text(value: Any, default: str = "") -> str:
    return str(value if value is not None else default).strip()


def extract_json(raw: str) -> list[dict]:
    cleaned = raw.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned, flags=re.I)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    start, end = cleaned.find("["), cleaned.rfind("]")
    if start < 0 or end < 0:
        raise ValueError("没有找到 JSON 数组。")
    data = json.loads(cleaned[start : end + 1])
    slides = []
    for i, item in enumerate(data, start=1):
        if not isinstance(item, dict):
            continue
        title = safe_text(item.get("title"), f"第 {i} 页")
        points = item.get("points", [])
        if isinstance(points, str):
            points = [points]
        points = [safe_text(p) for p in points if safe_text(p)]
        slide_type = safe_text(item.get("type"), "cards").lower()
        number = safe_text(item.get("number"), "")
        slides.append({"title": title, "points": points[:5], "type": slide_type, "number": number})
    return slides


def add_bg(slide, prs, theme):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill(bg, theme["bg"])
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.15), Inches(-0.55), Inches(2.8), Inches(2.8))
    fill(deco, theme["soft"], 0)
    deco2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.45), Inches(6.35), Inches(1.55), Inches(1.55))
    fill(deco2, theme["soft"], 0)


def add_header(slide, title, page_no, total, theme):
    text(slide, Inches(0.7), Inches(0.42), Inches(9.8), Inches(0.52), title, 25, theme["ink"], True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.14), Inches(1.05), Inches(0.08))
    fill(line, theme["accent"])
    text(slide, Inches(11.2), Inches(0.47), Inches(1.25), Inches(0.32), f"{page_no:02d}/{total:02d}", 10, theme["muted"], True, PP_ALIGN.RIGHT)


def add_footer(slide, theme):
    text(slide, Inches(0.72), Inches(6.86), Inches(3.2), Inches(0.22), "AI PPT Agent Pro", 8, theme["muted"])
    mark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.12), Inches(6.9), Inches(0.35), Inches(0.04))
    fill(mark, theme["accent"])


def cover(prs, topic, style, theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill(bg, theme["dark"])

    halo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.3), Inches(-0.9), Inches(5.2), Inches(5.2))
    fill(halo, theme["accent"], 12)
    halo2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.65), Inches(1.0), Inches(2.5), Inches(2.5))
    fill(halo2, theme["accent2"], 30)

    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(0.72), Inches(2.4), Inches(0.42))
    fill(tag, theme["accent"])
    text(s, Inches(0.78), Inches(0.72), Inches(2.4), Inches(0.42), "AI GENERATED DECK", 9, theme["white"], True, PP_ALIGN.CENTER)

    text(s, Inches(0.82), Inches(1.65), Inches(7.3), Inches(2.05), topic, 34, theme["white"], True, valign=MSO_ANCHOR.TOP)
    text(s, Inches(0.86), Inches(4.02), Inches(6.5), Inches(0.45), f"{style} · 精致版", 16, "DCEBFA")

    # 右侧抽象主视觉
    for i, (x, y, w, h, c) in enumerate([
        (8.2, 3.15, 3.7, 0.72, "FFFFFF"),
        (8.78, 4.08, 3.0, 0.58, "FFFFFF"),
        (8.35, 4.88, 2.55, 0.58, "FFFFFF"),
    ]):
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        fill(shape, c, 18 + i * 8)

    text(s, Inches(0.86), Inches(6.58), Inches(4.2), Inches(0.26), "Designed by AI PPT Agent Pro", 9, "BFD7EE")


def layout_cards(prs, data, page_no, total, theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs, theme)
    add_header(s, data["title"], page_no, total, theme)

    pts = data["points"] or ["本页内容待补充"]
    pts = pts[:4]
    positions = [(0.74, 1.62), (6.72, 1.62), (0.74, 4.02), (6.72, 4.02)]
    for i, p in enumerate(pts):
        x, y = positions[i]
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.55), Inches(1.78))
        fill(card, theme["card"])
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.09), Inches(1.78))
        fill(strip, theme["accent"])
        text(s, Inches(x + 0.35), Inches(y + 0.23), Inches(0.48), Inches(0.38), f"{i+1}", 13, theme["accent"], True, PP_ALIGN.CENTER)
        text(s, Inches(x + 0.98), Inches(y + 0.22), Inches(4.15), Inches(1.22), p, 17, theme["body"], valign=MSO_ANCHOR.TOP)
    add_footer(s, theme)


def layout_process(prs, data, page_no, total, theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs, theme)
    add_header(s, data["title"], page_no, total, theme)
    pts = (data["points"] or [])[:5]

    for i, p in enumerate(pts):
        x = 0.9 + i * 2.35
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.0), Inches(0.8), Inches(0.8))
        fill(circle, theme["accent"])
        text(s, Inches(x), Inches(2.0), Inches(0.8), Inches(0.8), f"{i+1}", 16, theme["white"], True, PP_ALIGN.CENTER)

        if i < len(pts) - 1:
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.86), Inches(2.36), Inches(1.15), Inches(0.06))
            fill(bar, theme["accent2"])

        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.25), Inches(3.1), Inches(1.7), Inches(1.65))
        fill(card, theme["card"])
        text(s, Inches(x - 0.05), Inches(3.3), Inches(1.3), Inches(1.05), p, 13, theme["body"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
    add_footer(s, theme)


def layout_split(prs, data, page_no, total, theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs, theme)
    add_header(s, data["title"], page_no, total, theme)
    pts = (data["points"] or [])[:4]

    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(1.62), Inches(5.0), Inches(4.75))
    fill(left, theme["dark"])
    text(s, Inches(1.2), Inches(2.0), Inches(4.2), Inches(0.68), data.get("number") or "01", 34, theme["white"], True)
    text(s, Inches(1.2), Inches(3.0), Inches(4.05), Inches(1.8), pts[0] if pts else "核心信息", 20, "DCEBFA", True, valign=MSO_ANCHOR.TOP)

    for i, p in enumerate(pts[1:] or ["补充说明"]):
        y = 1.78 + i * 1.35
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(y), Inches(0.55), Inches(0.43))
        fill(badge, theme["accent"])
        text(s, Inches(6.35), Inches(y), Inches(0.55), Inches(0.43), f"{i+1}", 10, theme["white"], True, PP_ALIGN.CENTER)
        text(s, Inches(7.08), Inches(y - 0.02), Inches(4.8), Inches(0.75), p, 17, theme["body"], valign=MSO_ANCHOR.TOP)
    add_footer(s, theme)


def layout_quote(prs, data, page_no, total, theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill(bg, theme["dark"])
    pts = data["points"] or ["感谢观看"]
    text(s, Inches(0.86), Inches(0.72), Inches(5.5), Inches(0.45), data["title"], 25, theme["white"], True)
    text(s, Inches(0.9), Inches(1.85), Inches(9.8), Inches(1.5), pts[0], 28, theme["white"], True, valign=MSO_ANCHOR.TOP)
    for i, p in enumerate(pts[1:4]):
        y = 4.05 + i * 0.72
        text(s, Inches(1.05), Inches(y), Inches(9.2), Inches(0.45), f"• {p}", 16, "DCEBFA")
    text(s, Inches(10.6), Inches(6.68), Inches(1.8), Inches(0.28), f"{page_no}/{total}", 9, "BFD7EE", align=PP_ALIGN.RIGHT)


def create_pptx(topic, style, slides, theme_name):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    theme = THEMES[theme_name]
    total = len(slides) + 1

    cover(prs, topic, style, theme)

    for page_no, item in enumerate(slides, start=2):
        t = item.get("type", "cards")
        title = item.get("title", "")
        if page_no == total or "总结" in title or "结论" in title:
            layout_quote(prs, item, page_no, total, theme)
        elif t in ["process", "流程", "timeline"]:
            layout_process(prs, item, page_no, total, theme)
        elif t in ["split", "重点", "highlight", "value"]:
            layout_split(prs, item, page_no, total, theme)
        else:
            layout_cards(prs, item, page_no, total, theme)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


st.markdown(
    """
<style>
/* 页面主体 */
.block-container {
    max-width: 1080px;
    padding-top: 1.6rem;
    padding-bottom: 2.5rem;
    padding-left: 1.2rem;
    padding-right: 1.2rem;
}

/* 卡片、指标与按钮 */
[data-testid="stMetric"] {
    background: #ffffff;
    border: 1px solid #ececec;
    padding: 12px;
    border-radius: 16px;
}

.stButton > button,
.stDownloadButton > button {
    min-height: 46px;
    border-radius: 14px;
    font-weight: 700;
    width: 100%;
}

div[data-testid="stVerticalBlockBorderWrapper"] {
    border-radius: 18px;
}

/* 手机端适配 */
@media (max-width: 768px) {
    .block-container {
        padding-top: 1rem;
        padding-left: 0.8rem;
        padding-right: 0.8rem;
        padding-bottom: 2rem;
    }

    h1 {
        font-size: 1.9rem !important;
        line-height: 1.2 !important;
    }

    p {
        line-height: 1.6;
    }

    /* 手机上让多列自动变成纵向 */
    div[data-testid="stHorizontalBlock"] {
        flex-direction: column !important;
        gap: 0.7rem !important;
    }

    div[data-testid="column"] {
        width: 100% !important;
        flex: 1 1 100% !important;
    }

    .stButton > button,
    .stDownloadButton > button {
        min-height: 50px;
        font-size: 1rem;
    }

    [data-testid="stMetric"] {
        margin-bottom: 0.4rem;
    }

    div[data-testid="stExpander"] {
        border-radius: 14px;
    }
}
</style>
""",
    unsafe_allow_html=True,
)

st.title("✨ AI PPT Agent Pro")
st.caption("手机和电脑都能使用：多版式模板 · 自动排版 · 一键下载")

with st.container(border=True):
    topic = st.text_input("PPT 主题", placeholder="例如：人工智能在矿物加工中的应用")
    page_count = st.slider("总页数", 5, 20, 8)

    option_col1, option_col2 = st.columns(2)
    with option_col1:
        style = st.selectbox(
            "内容风格",
            ["学术汇报", "商务汇报", "课堂展示", "简洁科普"],
        )
    with option_col2:
        theme_name = st.selectbox(
            "视觉主题",
            list(THEMES.keys()),
        )

    st.caption("手机端建议生成 6–10 页，速度更快，下载也更稳定。")

    generate = st.button(
        "✨ 生成精致版 PPT",
        type="primary",
        use_container_width=True,
    )

if generate:
    if not topic.strip():
        st.warning("请先输入 PPT 主题。")
    elif not api_key:
        st.error("没有读取到 DeepSeek API Key，请检查 .env 文件。")
    else:
        content_pages = page_count - 1
        prompt = f"""
请为主题“{topic}”设计一份 {style} 风格的 PowerPoint。

只返回 JSON 数组，不要解释，不要 Markdown，不要代码围栏。

格式：
[
  {{
    "title": "页面标题",
    "type": "cards/process/split/quote",
    "number": "可选的大数字或关键词",
    "points": ["要点1", "要点2", "要点3", "要点4"]
  }}
]

要求：
1. 生成 {content_pages} 页正文内容，不包含封面。
2. 必须混合使用 cards、process、split、quote 四种页面类型。
3. 最后一页必须是总结页，type 用 quote。
4. 每页 2 到 4 个要点。
5. 每个要点尽量控制在 28 个汉字以内。
6. 内容具体、逻辑紧凑，避免空话。
7. 全部使用中文。
"""

        with st.spinner("正在生成精致版 PPT..."):
            raw = ""
            try:
                client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com")
                response = client.chat.completions.create(
                    model="deepseek-v4-flash",
                    messages=[
                        {"role": "system", "content": "你是资深演示文稿策划师，必须严格返回可解析 JSON。"},
                        {"role": "user", "content": prompt},
                    ],
                )

                raw = response.choices[0].message.content or ""
                slides = extract_json(raw)
                ppt_file = create_pptx(topic, style, slides, theme_name)

                st.success("精致版 PPT 已生成。")
                st.info("手机用户：点击下方下载按钮后，可使用 WPS 或 Microsoft PowerPoint 打开。")
                m1, m2, m3 = st.columns(3)
                m1.metric("总页数", len(slides) + 1)
                m2.metric("内容风格", style)
                m3.metric("视觉主题", theme_name)

                with st.expander("查看页面结构"):
                    for i, slide in enumerate(slides, start=2):
                        st.markdown(f"**第 {i} 页｜{slide['title']}｜{slide.get('type','cards')}**")
                        for p in slide["points"]:
                            st.write(f"• {p}")

                safe = re.sub(r'[\\/:*?"<>|]+', "_", topic).strip() or "AI_PPT"
                st.download_button(
                    "⬇️ 下载精致版 PowerPoint",
                    data=ppt_file,
                    file_name=f"{safe}_精致版.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary",
                    use_container_width=True,
                )

            except json.JSONDecodeError:
                st.error("AI 返回格式不是有效 JSON，请重新生成。")
                with st.expander("查看原始返回"):
                    st.code(raw)
            except Exception as error:
                st.error(f"生成失败：{error}")
